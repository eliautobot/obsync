from __future__ import annotations

import json
from email.message import EmailMessage
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from obsync.extractors import extract_document


def test_text_encoding_and_truncation(tmp_path: Path) -> None:
    path = tmp_path / "notes.txt"
    path.write_bytes("Café and project notes".encode("cp1252"))
    extracted = extract_document(path, max_chars=10)
    assert extracted.extractor == "text"
    assert extracted.truncated
    assert "Content truncated" in extracted.text


def test_json_is_pretty_printed(tmp_path: Path) -> None:
    path = tmp_path / "record.json"
    path.write_text(json.dumps({"name": "Acme", "status": "active"}), encoding="utf-8")
    extracted = extract_document(path)
    assert extracted.extractor == "json"
    assert '\n  "name": "Acme"' in extracted.text


def test_csv_rows_become_readable_lines(tmp_path: Path) -> None:
    path = tmp_path / "items.csv"
    path.write_text("name,total\nPipe,12.50\nValve,8.00\n", encoding="utf-8")
    extracted = extract_document(path)
    assert "name | total" in extracted.text
    assert "Pipe | 12.50" in extracted.text


def test_docx_paragraphs_and_tables(tmp_path: Path) -> None:
    path = tmp_path / "report.docx"
    document = Document()
    document.add_paragraph("Project report")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Status"
    table.rows[0].cells[1].text = "Complete"
    document.save(path)
    extracted = extract_document(path)
    assert "Project report" in extracted.text
    assert "Status | Complete" in extracted.text


def test_xlsx_values_are_extracted(tmp_path: Path) -> None:
    path = tmp_path / "invoices.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Invoices"
    sheet.append(["Number", "Total"])
    sheet.append([1001, 99.5])
    workbook.save(path)
    extracted = extract_document(path)
    assert "## Sheet: Invoices" in extracted.text
    assert "1001 | 99.5" in extracted.text


def test_pptx_slide_text_is_extracted(tmp_path: Path) -> None:
    path = tmp_path / "briefing.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly briefing"
    slide.placeholders[1].text = "Operations are on schedule"
    presentation.save(path)
    extracted = extract_document(path)
    assert "Quarterly briefing" in extracted.text
    assert "Operations are on schedule" in extracted.text


def test_unknown_binary_gets_metadata_fallback(tmp_path: Path) -> None:
    path = tmp_path / "archive.bin"
    path.write_bytes(b"\x00\x01\x02")
    extracted = extract_document(path)
    assert extracted.extractor == "metadata"
    assert extracted.text == ""
    assert "No text extractor" in extracted.warning


def test_html_scripts_are_removed(tmp_path: Path) -> None:
    path = tmp_path / "page.html"
    path.write_text(
        "<html><style>.hidden{}</style><script>secret()</script><body><h1>Visible</h1></body></html>",
        encoding="utf-8",
    )
    extracted = extract_document(path)
    assert "Visible" in extracted.text
    assert "secret" not in extracted.text


def test_eml_headers_and_body_are_extracted(tmp_path: Path) -> None:
    message = EmailMessage()
    message["Subject"] = "Project update"
    message["From"] = "sender@example.com"
    message["To"] = "team@example.com"
    message.set_content("The project is on schedule.")
    path = tmp_path / "message.eml"
    path.write_bytes(message.as_bytes())
    extracted = extract_document(path)
    assert "Subject: Project update" in extracted.text
    assert "The project is on schedule" in extracted.text
