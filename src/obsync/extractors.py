from __future__ import annotations

import csv
import json
import mimetypes
from collections.abc import Callable
from dataclasses import dataclass
from email import policy
from email.parser import BytesParser
from pathlib import Path

from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".rst",
    ".log",
    ".csv",
    ".tsv",
    ".json",
    ".jsonl",
    ".yaml",
    ".yml",
    ".toml",
    ".xml",
    ".ini",
    ".cfg",
    ".conf",
    ".py",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".css",
    ".scss",
    ".sql",
    ".sh",
    ".ps1",
    ".bat",
    ".go",
    ".rs",
    ".java",
    ".cs",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp"}


@dataclass(slots=True)
class ExtractedDocument:
    text: str
    mime_type: str
    extractor: str
    truncated: bool = False
    warning: str = ""


def detect_mime(path: Path) -> str:
    return mimetypes.guess_type(path.name)[0] or "application/octet-stream"


def _decode_bytes(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-16", "cp1252", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_text(path: Path) -> str:
    return _decode_bytes(path.read_bytes())


def _extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    parts = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"## Page {index}\n\n{text.strip()}")
    return "\n\n".join(parts)


def _extract_docx(path: Path) -> str:
    document = Document(str(path))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table_index, table in enumerate(document.tables, start=1):
        parts.append(f"\nTable {table_index}")
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    workbook = load_workbook(path, read_only=True, data_only=True)
    output: list[str] = []
    cell_count = 0
    max_cells = 50_000
    try:
        for sheet in workbook.worksheets:
            output.append(f"## Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(values):
                    output.append(" | ".join(values))
                cell_count += len(values)
                if cell_count >= max_cells:
                    output.append("[Spreadsheet extraction stopped after 50,000 cells]")
                    return "\n".join(output)
    finally:
        workbook.close()
    return "\n".join(output)


def _extract_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    output: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        output.append(f"## Slide {index}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                output.append(shape.text.strip())
    return "\n".join(output)


def _extract_html(path: Path) -> str:
    soup = BeautifulSoup(_decode_bytes(path.read_bytes()), "html.parser")
    for element in soup(["script", "style", "noscript"]):
        element.decompose()
    return "\n".join(line.strip() for line in soup.get_text("\n").splitlines() if line.strip())


def _extract_eml(path: Path) -> str:
    message = BytesParser(policy=policy.default).parsebytes(path.read_bytes())
    headers = [
        f"Subject: {message.get('subject', '')}",
        f"From: {message.get('from', '')}",
        f"To: {message.get('to', '')}",
        f"Date: {message.get('date', '')}",
    ]
    body = message.get_body(preferencelist=("plain", "html")) if message.is_multipart() else message
    content = ""
    if body:
        try:
            content = body.get_content()
        except (LookupError, UnicodeDecodeError):
            payload = body.get_payload(decode=True) or b""
            content = _decode_bytes(payload)
    if body and body.get_content_type() == "text/html":
        content = BeautifulSoup(content, "html.parser").get_text("\n")
    return "\n".join(headers) + "\n\n" + str(content)


def _extract_csv(path: Path) -> str:
    text = _decode_bytes(path.read_bytes())
    sample = text[:4096]
    try:
        dialect = csv.Sniffer().sniff(sample)
    except csv.Error:
        return text
    rows = []
    for row in csv.reader(text.splitlines(), dialect):
        rows.append(" | ".join(row))
    return "\n".join(rows)


def _extract_json(path: Path) -> str:
    text = _decode_bytes(path.read_bytes())
    try:
        return json.dumps(json.loads(text), ensure_ascii=False, indent=2)
    except json.JSONDecodeError:
        return text


def _extract_image(path: Path) -> tuple[str, str]:
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        return "", "OCR dependencies are not installed"
    try:
        with Image.open(path) as image:
            return pytesseract.image_to_string(image), ""
    except (OSError, RuntimeError) as exc:
        return "", f"OCR failed: {exc}"


def extract_document(path: Path, max_chars: int = 200_000) -> ExtractedDocument:
    extension = path.suffix.lower()
    mime_type = detect_mime(path)
    warning = ""

    extractor: str
    function: Callable[[Path], str] | None
    if extension == ".pdf":
        extractor, function = "pypdf", _extract_pdf
    elif extension == ".docx":
        extractor, function = "python-docx", _extract_docx
    elif extension in {".xlsx", ".xlsm", ".xltx"}:
        extractor, function = "openpyxl", _extract_xlsx
    elif extension == ".pptx":
        extractor, function = "python-pptx", _extract_pptx
    elif extension in {".html", ".htm"}:
        extractor, function = "beautifulsoup", _extract_html
    elif extension == ".eml":
        extractor, function = "email", _extract_eml
    elif extension in {".csv", ".tsv"}:
        extractor, function = "csv", _extract_csv
    elif extension in {".json", ".jsonl"}:
        extractor, function = "json", _extract_json
    elif extension in IMAGE_EXTENSIONS:
        extractor, function = "tesseract", None
    elif extension in TEXT_EXTENSIONS or mime_type.startswith("text/"):
        extractor, function = "text", _extract_text
    else:
        extractor, function = "metadata", None

    if extension in IMAGE_EXTENSIONS:
        text, warning = _extract_image(path)
    elif function:
        text = function(path)
    else:
        text = ""
        warning = "No text extractor is available for this file type; metadata was synced."

    text = text.replace("\x00", "").strip()
    truncated = len(text) > max_chars
    if truncated:
        text = text[:max_chars].rstrip() + "\n\n[Content truncated by Obsync]"
    return ExtractedDocument(
        text=text,
        mime_type=mime_type,
        extractor=extractor,
        truncated=truncated,
        warning=warning,
    )
